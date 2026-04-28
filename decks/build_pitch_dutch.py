"""
Build the Hey Otis pitch deck — Dutch incubator edition (.pptx).

Tuned for Antler NL, Rockstart Health, TechLeap, ACE Incubator.
European-first framing, capital-efficient burn, AI Act + WBSO mentions,
two-founder team (Samantha + husband co-founder).

Run:  python3 decks/build_pitch_dutch.py
Out:  decks/HeyOtis_Pitch_Dutch.pptx

The last 4 slides are "Why this incubator" variants — keep one,
delete the other three before submitting per application.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Brand
INK = RGBColor(0x21, 0x1E, 0x28)
CREAM = RGBColor(0xF7, 0xF5, 0xFD)
LIME = RGBColor(0x96, 0xD3, 0x5F)
LIME_LIGHT = RGBColor(0xC8, 0xEC, 0xB0)
PERIWINKLE = RGBColor(0x92, 0xA6, 0xF4)
ORANGE = RGBColor(0xF6, 0x77, 0x00)
PURPLE = RGBColor(0xBD, 0x57, 0xF2)
TEAL = RGBColor(0x4E, 0xA9, 0x89)
SUB = RGBColor(0x80, 0x79, 0x8C)
BODY = RGBColor(0x3A, 0x36, 0x30)
EYEBROW_GREEN = RGBColor(0x4A, 0x7A, 0x23)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    bg_shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg_shape.line.fill.background()
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    return s


def add_text(slide, text, left, top, width, height,
             font="Inter", size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.3):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    # Allow newlines in input
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tx


def eyebrow(slide, text, left, top, width=Inches(8), color=EYEBROW_GREEN):
    add_text(slide, text.upper(), left, top, width, Inches(0.4),
             size=12, bold=True, color=color)


def add_card(slide, left, top, width, height, fill=WHITE,
             stroke=RGBColor(0xDE, 0xDD, 0xE8), thick=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.06
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = stroke
    shape.line.width = Pt(2 if thick else 0.75)
    shape.shadow.inherit = False
    return shape


def section_title(slide, eyebrow_text, title, eyebrow_color=EYEBROW_GREEN,
                  title_size=36, title_top=Inches(1.4)):
    eyebrow(slide, eyebrow_text, Inches(0.8), Inches(0.8), color=eyebrow_color)
    add_text(slide, title, Inches(0.8), title_top, Inches(12), Inches(1.4),
             size=title_size, bold=True, line_spacing=1.15)


# ────────────────────────────────────────────────────────────────────────
# 01 — COVER
# ────────────────────────────────────────────────────────────────────────
s = add_slide(LIME_LIGHT)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.5), SW, Inches(3))
band.line.fill.background()
band.fill.solid(); band.fill.fore_color.rgb = LIME
add_text(s, "hey Otis", Inches(0.8), Inches(0.6), Inches(6), Inches(1),
         size=40, bold=True, color=INK)
add_text(s, "From rupture to repair.", Inches(0.8), Inches(2.4), Inches(11), Inches(2),
         size=72, bold=True, color=INK, line_spacing=1.0)
add_text(s, "A private AI guide for couples in conflict.",
         Inches(0.8), Inches(4.7), Inches(11), Inches(1),
         size=22, color=INK, line_spacing=1.3)
add_text(s, "Samantha Ahhee  ·  Founder, Hey Otis  ·  Amsterdam  ·  2026",
         Inches(0.8), Inches(6.6), Inches(11), Inches(0.5),
         size=13, bold=True, color=INK)


# ────────────────────────────────────────────────────────────────────────
# 02 — THE PROBLEM
# ────────────────────────────────────────────────────────────────────────
s = add_slide(WHITE)
section_title(s, "The problem", "Conflict isn't the problem. Silence is.")
add_text(s,
         "69% of relationship conflicts never fully resolve.\n"
         "Couples wait six years before seeking help.\n"
         "Most fights happen at 11pm — nine days from the next therapy "
         "session, if there is one.",
         Inches(0.8), Inches(2.8), Inches(11.5), Inches(2.5),
         size=22, color=BODY, line_spacing=1.5)
add_text(s,
         "When my own relationship hit that wall, I had nothing to reach for "
         "that wasn't a self-help book or a 3am text to a friend.",
         Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5),
         size=18, color=SUB, line_spacing=1.5)


# ────────────────────────────────────────────────────────────────────────
# 03 — WHY NOW
# ────────────────────────────────────────────────────────────────────────
s = add_slide(CREAM)
section_title(s, "Why now", "Three things shifted in the last 24 months.")
items = [
    ("Models are ready.", "LLMs crossed the threshold for emotional nuance without giving harmful advice."),
    ("Regulation is here.", "The EU AI Act gives us a clear frame to build trustworthy mental-health AI on."),
    ("Demand broke open.", "Dutch GGZ relatietherapie waitlists hit 6+ months. The gap is now structural."),
]
top = Inches(2.9)
for i, (h, b) in enumerate(items):
    y = top + Inches(0.05) + Inches(1.35) * i
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.15),
                             Inches(0.25), Inches(0.25))
    dot.fill.solid(); dot.fill.fore_color.rgb = LIME
    dot.line.fill.background()
    add_text(s, h, Inches(1.3), y, Inches(11), Inches(0.6),
             size=22, bold=True)
    add_text(s, b, Inches(1.3), y + Inches(0.6), Inches(11), Inches(0.7),
             size=15, color=SUB, line_spacing=1.5)


# ────────────────────────────────────────────────────────────────────────
# 04 — THE PRODUCT
# ────────────────────────────────────────────────────────────────────────
s = add_slide(WHITE)
eyebrow(s, "The product", Inches(0.8), Inches(0.8))
add_text(s, "A private space to use during conflict.",
         Inches(0.8), Inches(1.4), Inches(7.5), Inches(1.6),
         size=32, bold=True, line_spacing=1.15)

bullets = [
    "Open it after a fight. Walk through four steps. Arrive at a conversation.",
    "Solo or synced with your partner. Voice or text input.",
    "End-to-end private. We never train on user data. GDPR-native.",
    "Available on iOS, Android, and web.",
]
ty = Inches(3.4)
for b in bullets:
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), ty + Inches(0.13),
                             Inches(0.18), Inches(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = LIME
    dot.line.fill.background()
    add_text(s, b, Inches(1.15), ty, Inches(6.5), Inches(0.6),
             size=15, color=BODY, line_spacing=1.5)
    ty += Inches(0.65)

# Optional product image right side
try:
    s.shapes.add_picture(
        "/Users/samanthaahhee/tether/website/public/hero-image.png",
        Inches(8.0), Inches(1.5), width=Inches(4.7),
    )
except FileNotFoundError:
    add_card(s, Inches(8.0), Inches(1.5), Inches(4.7), Inches(5), fill=CREAM)
    add_text(s, "Product screen", Inches(8.0), Inches(3.8), Inches(4.7), Inches(0.5),
             size=14, color=SUB, align=PP_ALIGN.CENTER)


# ────────────────────────────────────────────────────────────────────────
# 05 — HOW IT WORKS — 4 step row
# ────────────────────────────────────────────────────────────────────────
s = add_slide(CREAM)
section_title(s, "How it works", "Four steps. One conversation at a time.")

steps = [
    (1, TEAL, "Vent", "Say what you're feeling. Privately."),
    (2, PERIWINKLE, "Understand", "Find the unmet need underneath."),
    (3, ORANGE, "Prepare", "Frame what you want to say."),
    (4, PURPLE, "Nurture", "Walk through the actual repair."),
]
cw = Inches(2.85)
gap = Inches(0.2)
total = cw * 4 + gap * 3
start_left = (SW - total) // 2
top = Inches(3.1)
for i, (n, col, name, body) in enumerate(steps):
    left = start_left + (cw + gap) * i
    add_card(s, left, top, cw, Inches(2.8))
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, left + (cw - Inches(0.7)) // 2,
                                top + Inches(0.4), Inches(0.7), Inches(0.7))
    circle.fill.solid(); circle.fill.fore_color.rgb = col
    circle.line.fill.background()
    add_text(s, str(n), left, top + Inches(0.46), cw, Inches(0.6),
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, line_spacing=1.0)
    add_text(s, name, left, top + Inches(1.3), cw, Inches(0.5),
             size=20, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, body, left + Inches(0.2), top + Inches(1.9),
             cw - Inches(0.4), Inches(0.9), size=12, color=SUB,
             align=PP_ALIGN.CENTER, line_spacing=1.5)

add_text(s,
         "Each step maps to peer-reviewed couples research: emotion regulation, "
         "attachment, NVC, Gottman softened start-up.",
         Inches(0.8), Inches(6.4), Inches(11.5), Inches(1),
         size=13, color=SUB, align=PP_ALIGN.CENTER, line_spacing=1.5)


# ────────────────────────────────────────────────────────────────────────
# 06 — WEDGE & DEFENSIBILITY
# ────────────────────────────────────────────────────────────────────────
s = add_slide(WHITE)
section_title(s, "Wedge & defensibility", "Why ChatGPT won't eat this. Why BetterHelp can't.")

cards = [
    ("Frameworks library",
     "Validated by clinical advisor. 18 months to build, hard to clone."),
    ("Memory flywheel",
     "Each session compounds. Generic chatbots forget you between chats."),
    ("Brand trust",
     "EU-native, private by default. The category punishes any breach."),
]
cw = Inches(3.85)
gap = Inches(0.25)
start_left = (SW - (cw * 3 + gap * 2)) // 2
top = Inches(3.0)
for i, (title, body) in enumerate(cards):
    left = start_left + (cw + gap) * i
    add_card(s, left, top, cw, Inches(2.6))
    add_text(s, title, left + Inches(0.3), top + Inches(0.4),
             cw - Inches(0.6), Inches(0.7), size=18, bold=True)
    add_text(s, body, left + Inches(0.3), top + Inches(1.2),
             cw - Inches(0.6), Inches(1.5), size=13, color=SUB, line_spacing=1.5)

add_text(s,
         "We are explicitly not a medical device. The clinical advisor on "
         "the team validates that boundary monthly.",
         Inches(0.8), Inches(6.4), Inches(11.5), Inches(1),
         size=13, color=SUB, align=PP_ALIGN.CENTER)


# ────────────────────────────────────────────────────────────────────────
# 07 — MARKET
# ────────────────────────────────────────────────────────────────────────
s = add_slide(CREAM)
section_title(s, "The market", "European-first. Bottoms-up.")

rows = [
    ("47M", "long-term couples in the EU"),
    ("38M", "are not in any form of therapy"),
    ("0.5%", "= 190k subs in 5 years × €15/mo = €34M ARR"),
]
ry = Inches(2.9)
for n, lab in rows:
    add_card(s, Inches(0.8), ry, Inches(11.7), Inches(1.0))
    add_text(s, n, Inches(1.1), ry + Inches(0.18), Inches(2.5), Inches(0.7),
             size=30, bold=True, color=TEAL)
    add_text(s, lab, Inches(3.7), ry + Inches(0.28), Inches(8.5), Inches(0.5),
             size=15, color=BODY)
    ry += Inches(1.15)

add_text(s,
         "Comparables: Replika (10M users, $80M ARR) · Paired UK ($9M Series A) · "
         "Headspace (€600M valuation, mental-wellness adjacency).",
         Inches(0.8), Inches(6.4), Inches(11.5), Inches(1),
         size=12, color=SUB, align=PP_ALIGN.CENTER, line_spacing=1.5)
add_text(s,
         "We start in NL + UK + DACH. US is Year 3, not Year 1.",
         Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.4),
         size=13, bold=True, color=EYEBROW_GREEN, align=PP_ALIGN.CENTER)


# ────────────────────────────────────────────────────────────────────────
# 08 — BUSINESS MODEL
# ────────────────────────────────────────────────────────────────────────
s = add_slide(WHITE)
section_title(s, "Business model", "Subscription. Couple-tier is the unlock.")

prices = [
    ("Free", "€0", "One guided repair per month\nbasic insights", False),
    ("Plus (solo)", "€14.99/mo", "Unlimited repairs, partner sync,\nfull assessment library", True),
    ("Couples (both)", "€22/mo", "Both partners synced,\nshared history, deeper insight", False),
]
pw = Inches(3.85)
gap = Inches(0.25)
total = pw * 3 + gap * 2
start_left = (SW - total) // 2
top = Inches(2.9)
for i, (tier, price, body, feat) in enumerate(prices):
    left = start_left + (pw + gap) * i
    add_card(s, left, top, pw, Inches(2.6),
             stroke=LIME if feat else RGBColor(0xDE, 0xDD, 0xE8),
             thick=feat)
    add_text(s, tier.upper(), left, top + Inches(0.3),
             pw, Inches(0.4), size=12, bold=True, color=SUB,
             align=PP_ALIGN.CENTER)
    add_text(s, price, left, top + Inches(0.85),
             pw, Inches(0.8), size=30, bold=True, color=INK,
             align=PP_ALIGN.CENTER)
    add_text(s, body, left + Inches(0.3), top + Inches(1.7),
             pw - Inches(0.6), Inches(1), size=12, color=SUB,
             align=PP_ALIGN.CENTER, line_spacing=1.5)

add_text(s,
         "Assumptions: blended CAC €18  ·  14-month tenure  ·  78% gross margin  ·  LTV €185.",
         Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.5),
         size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
add_text(s,
         "Couple subscription = built-in viral loop. One signup brings two users.",
         Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.5),
         size=13, color=SUB, align=PP_ALIGN.CENTER)


# ────────────────────────────────────────────────────────────────────────
# 09 — YEAR 1 GTM
# ────────────────────────────────────────────────────────────────────────
s = add_slide(CREAM)
section_title(s, "Year 1 GTM", "Three channels. Honest CPA targets.")

channels = [
    ("Reddit + Substack organic",
     "Relationship subs, attachment-style writers. Already running.",
     "CPA target: €4"),
    ("Couples-therapist referrals",
     "Clinicians give Hey Otis to clients between sessions.",
     "CPA target: €8"),
    ("Instagram Reels + creators",
     "Short-form clinical-content partnerships. Brand-driven.",
     "CPA target: €25"),
]
top = Inches(2.9)
for i, (title, body, cpa) in enumerate(channels):
    y = top + Inches(1.15) * i
    add_card(s, Inches(0.8), y, Inches(11.7), Inches(1.0))
    add_text(s, title, Inches(1.1), y + Inches(0.2), Inches(7), Inches(0.5),
             size=18, bold=True)
    add_text(s, body, Inches(1.1), y + Inches(0.55), Inches(8), Inches(0.4),
             size=12, color=SUB)
    add_text(s, cpa, Inches(9.5), y + Inches(0.3), Inches(2.8), Inches(0.5),
             size=15, bold=True, color=EYEBROW_GREEN, align=PP_ALIGN.RIGHT)

add_text(s,
         "Channel 1 is already validated by the heyotis.app/research survey campaign.",
         Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
         size=13, color=SUB, align=PP_ALIGN.CENTER)


# ────────────────────────────────────────────────────────────────────────
# 10 — TRACTION
# ────────────────────────────────────────────────────────────────────────
s = add_slide(WHITE)
section_title(s, "Traction", "Where we are today.")

big = [
    ("[X]", "Waitlist signups", "heyotis.app, organic"),
    ("[X]", "Survey responses", "primary research, real couples"),
    ("78%", "Said they'd try it", "from in-survey validation"),
]
cw = Inches(3.85)
gap = Inches(0.25)
start_left = (SW - (cw * 3 + gap * 2)) // 2
top = Inches(2.9)
for i, (n, lab, sub) in enumerate(big):
    left = start_left + (cw + gap) * i
    add_card(s, left, top, cw, Inches(2.5))
    add_text(s, n, left, top + Inches(0.45),
             cw, Inches(1), size=52, bold=True, color=INK,
             align=PP_ALIGN.CENTER, line_spacing=1.0)
    add_text(s, lab.upper(), left, top + Inches(1.55),
             cw, Inches(0.4), size=12, bold=True, color=INK,
             align=PP_ALIGN.CENTER)
    add_text(s, sub, left + Inches(0.3), top + Inches(1.95),
             cw - Inches(0.6), Inches(0.6), size=11, color=SUB,
             align=PP_ALIGN.CENTER)

add_text(s,
         "Brand, marketing site, Supabase backend, payments, app prototype — all live.\n"
         "Beta launching Q3 2026 in EN-speaking markets. Apple + Google submission ready.",
         Inches(0.8), Inches(6.0), Inches(11.5), Inches(1.2),
         size=13, color=SUB, align=PP_ALIGN.CENTER, line_spacing=1.6)


# ────────────────────────────────────────────────────────────────────────
# 11 — TEAM
# ────────────────────────────────────────────────────────────────────────
s = add_slide(CREAM)
section_title(s, "Team", "Two founders. Lived experience + commercial pedigree.")

team = [
    ("Samantha Ahhee", "Founder + CEO",
     "Designer + builder. Previously [your prior role]. Built the Hey Otis "
     "product, brand, and primary-research system end-to-end. Lives the problem."),
    ("[Husband's name]", "Co-founder + GTM",
     "[N] years bringing software products to market. Previously [his prior "
     "role at company]. Owns commercial strategy, partnerships, and scale."),
    ("Hiring next", "Clinical lead",
     "Licensed couples therapist to validate frameworks, prompts, safety "
     "guardrails. In conversation with two candidates."),
]
cw = Inches(3.85)
gap = Inches(0.25)
start_left = (SW - (cw * 3 + gap * 2)) // 2
top = Inches(2.7)
for i, (name, role, body) in enumerate(team):
    left = start_left + (cw + gap) * i
    add_card(s, left, top, cw, Inches(3.7))
    add_text(s, name, left + Inches(0.3), top + Inches(0.35),
             cw - Inches(0.6), Inches(0.7), size=18, bold=True)
    add_text(s, role.upper(), left + Inches(0.3), top + Inches(1.0),
             cw - Inches(0.6), Inches(0.4), size=11, bold=True, color=EYEBROW_GREEN)
    add_text(s, body, left + Inches(0.3), top + Inches(1.5),
             cw - Inches(0.6), Inches(2), size=12, color=BODY, line_spacing=1.5)


# ────────────────────────────────────────────────────────────────────────
# 12 — ASK
# ────────────────────────────────────────────────────────────────────────
s = add_slide(LIME_LIGHT)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.5), SW, Inches(3))
band.line.fill.background(); band.fill.solid(); band.fill.fore_color.rgb = LIME
eyebrow(s, "The ask", Inches(0.8), Inches(0.9))
add_text(s, "Raising [€XXX]K pre-seed",
         Inches(0.5), Inches(1.7), Inches(12.3), Inches(2),
         size=58, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.1)
add_text(s,
         "24-month runway. To launch publicly, hit 10,000 paying users,\n"
         "validate Couples-tier LTV, and reach Seed milestones.",
         Inches(0.5), Inches(3.6), Inches(12.3), Inches(1.5),
         size=20, align=PP_ALIGN.CENTER, line_spacing=1.5)

# Non-dilutive callout
add_text(s,
         "Plan to stack with WBSO R&D credit + RVO Innovation Credit\n"
         "(targeting +6 months runway non-dilutively)",
         Inches(0.5), Inches(5.0), Inches(12.3), Inches(1),
         size=14, bold=True, color=EYEBROW_GREEN, align=PP_ALIGN.CENTER, line_spacing=1.5)

add_text(s, "samantha.ahhee@gmail.com  ·  heyotis.app",
         Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
         size=18, bold=True, align=PP_ALIGN.CENTER)


# ────────────────────────────────────────────────────────────────────────
# 13–16 — APPENDIX: Why this incubator (variants — keep one, delete rest)
# ────────────────────────────────────────────────────────────────────────
def why_slide(name, body, color=PERIWINKLE):
    s = add_slide(WHITE)
    eyebrow(s, "Appendix · Why this incubator (delete the others)",
            Inches(0.8), Inches(0.6), color=SUB)
    add_text(s, f"Why {name}", Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.2),
             size=42, bold=True, color=color, line_spacing=1.1)
    add_text(s, body, Inches(0.8), Inches(2.9), Inches(11.5), Inches(4),
             size=18, color=BODY, line_spacing=1.6)


why_slide("Antler NL",
    "Antler's residency model and consumer-health portfolio fit Hey Otis "
    "precisely. We come in with the product, brand, and primary research "
    "already in market. We use the residency to (a) validate clinical "
    "co-investor relationships, (b) sharpen the EU-first GTM motion, and "
    "(c) close the clinical-lead hire. The category we're in — couples "
    "mental wellness — is brand-and-trust before it's anything else, and "
    "Antler portfolio peers have solved exactly that distribution problem.",
    color=PERIWINKLE)

why_slide("Rockstart Health",
    "Rockstart Health's clinician network is the unlock for our second "
    "GTM channel: therapist referrals. Hey Otis only works as a "
    "therapist-recommended tool if therapists know it exists. Your "
    "portfolio of companies bridging digital and in-person care gives us "
    "the peer group that has solved exactly this distribution problem. "
    "We'd also leverage your brand to accelerate insurance-reimbursement "
    "conversations in Germany, where the pathway exists.",
    color=TEAL)

why_slide("TechLeap.nl",
    "TechLeap is the validation layer that opens the next two doors: the "
    "RVO Innovation Credit (for clinical-validation work) and the EIC "
    "Accelerator at Horizon Europe (for EU-wide expansion). Both are "
    "non-dilutive and well-suited to our 18-month roadmap. A TechLeap "
    "stamp also accelerates the brand-trust conversation that drives our "
    "B2C unit economics. We'd use the TechLeap network to bridge into "
    "both grant pipelines and Rise programme alumni for hiring.",
    color=ORANGE)

why_slide("ACE Incubator",
    "ACE's link to the University of Amsterdam — and the clinical "
    "psychology department specifically — gives Hey Otis the academic "
    "validation layer that consumer mental-health products usually have "
    "to build years in. A research collaboration with UvA on "
    "conflict-resolution outcomes would be the strongest defensibility "
    "moat we could build at this stage. ACE's local network also "
    "shortens the hiring loop for the clinical lead role.",
    color=PURPLE)


# ── Save ──
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "HeyOtis_Pitch_Dutch.pptx")
prs.save(out)
print(f"\n✓ Wrote {out}")
print(f"  Slides: {len(prs.slides)}")
print(f"\n  Main deck = slides 1–12")
print(f"  Appendix = slides 13–16 (Why-this-incubator variants — keep one, delete the others)")
