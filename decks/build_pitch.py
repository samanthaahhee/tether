"""
Build the Hey Otis pitch deck (.pptx).

Run:  python3 decks/build_pitch.py
Out:  decks/HeyOtis_Pitch.pptx

Upload to Google Slides:
  1. drive.google.com -> New -> File upload -> select the .pptx
  2. Right-click the uploaded file -> Open with -> Google Slides
  Google converts it; you can edit everything natively from there.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Brand tokens ────────────────────────────────────────────────────────
INK = RGBColor(0x21, 0x1E, 0x28)
CREAM = RGBColor(0xF7, 0xF5, 0xFD)
LIME = RGBColor(0x96, 0xD3, 0x5F)
LIME_LIGHT = RGBColor(0xC8, 0xEC, 0xB0)
PERIWINKLE = RGBColor(0x92, 0xA6, 0xF4)
ORANGE = RGBColor(0xF6, 0x77, 0x00)
PURPLE = RGBColor(0xBD, 0x57, 0xF2)
TEAL = RGBColor(0x4E, 0xA9, 0x89)
SUB = RGBColor(0x80, 0x79, 0x8C)
BODY = RGBColor(0x3A, 0x36, 0x30)
EYEBROW_GREEN = RGBColor(0x4A, 0x7A, 0x23)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9 widescreen, matches Google Slides
prs.slide_height = Inches(7.5)

SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    bg_shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg_shape.line.fill.background()
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    return s


def add_text(slide, text, left, top, width, height,
             font="Inter", size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.3):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tx


def eyebrow(slide, text, left, top, width=Inches(8), color=EYEBROW_GREEN):
    add_text(slide, text.upper(), left, top, width, Inches(0.4),
             size=12, bold=True, color=color)


def add_card(slide, left, top, width, height, fill=WHITE, stroke=RGBColor(0xDE, 0xDD, 0xE8), radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.adjustments[0] = 0.08 if radius else 0
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = stroke
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


# ──────────────────────────────────────────────────────────────────────────
# 01 — COVER
# ──────────────────────────────────────────────────────────────────────────
s = add_slide(LIME_LIGHT)
# top gradient feel via a second band
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.5), SW, Inches(3))
band.line.fill.background()
band.fill.solid()
band.fill.fore_color.rgb = LIME
add_text(s, "hey Otis", Inches(0.8), Inches(0.6), Inches(6), Inches(1),
         size=44, bold=True, color=INK)
add_text(s, "From rupture to repair.", Inches(0.8), Inches(2.4), Inches(11), Inches(2),
         size=80, bold=True, color=INK, line_spacing=1.0)
add_text(s, "A private guide for couples in conflict.",
         Inches(0.8), Inches(4.7), Inches(11), Inches(1),
         size=24, color=INK, line_spacing=1.3)
add_text(s, "Samantha Ahhee  ·  Founder, Hey Otis  ·  Amsterdam, 2026",
         Inches(0.8), Inches(6.6), Inches(11), Inches(0.5),
         size=13, bold=True, color=INK)


# ──────────────────────────────────────────────────────────────────────────
# 02 — HOOK
# ──────────────────────────────────────────────────────────────────────────
s = add_slide(WHITE)
eyebrow(s, "The hook", Inches(0.8), Inches(0.8))
add_text(s, "Every couple argues.\nThe ones that last learn how to repair.",
         Inches(0.8), Inches(1.6), Inches(11.5), Inches(2.5),
         size=46, bold=True, line_spacing=1.15)
add_text(s,
         "“Conflict in relationships is inevitable. It's not a sign that something is broken, "
         "it's a signal that something needs attention.”",
         Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.4),
         size=22, color=BODY, line_spacing=1.5)
add_text(s, "— Esther Perel", Inches(0.8), Inches(6.4), Inches(11), Inches(0.5),
         size=14, bold=True, color=EYEBROW_GREEN)


# ──────────────────────────────────────────────────────────────────────────
# 03 — PROBLEM (4 stat cards)
# ──────────────────────────────────────────────────────────────────────────
s = add_slide(CREAM)
eyebrow(s, "The problem", Inches(0.8), Inches(0.8))
add_text(s, "Conflict isn't the problem. Silence is.",
         Inches(0.8), Inches(1.4), Inches(12), Inches(1),
         size=40, bold=True)

stats = [
    ("69%", TEAL, "of conflicts are perpetual,\nthey never fully resolve"),
    ("6 yrs", PERIWINKLE, "couples wait before\nseeking any help"),
    ("96%", ORANGE, "of conversations are\ndetermined by how they start"),
    ("5:1", PURPLE, "positive to negative,\nthe magic ratio"),
]
card_w = Inches(2.85)
gap = Inches(0.2)
total_w = card_w * 4 + gap * 3
start_left = (SW - total_w) // 2
top = Inches(3.0)
for i, (n, col, lab) in enumerate(stats):
    left = start_left + (card_w + gap) * i
    add_card(s, left, top, card_w, Inches(2.5))
    add_text(s, n, left, top + Inches(0.3), card_w, Inches(1),
             size=44, bold=True, color=col, align=PP_ALIGN.CENTER, line_spacing=1.0)
    add_text(s, lab, left + Inches(0.2), top + Inches(1.5), card_w - Inches(0.4), Inches(1),
             size=12, color=SUB, align=PP_ALIGN.CENTER, line_spacing=1.4)

add_text(s,
         "Most couples don't break up because they fight. They break up because "
         "they stop trying to understand each other.",
         Inches(0.8), Inches(6.2), Inches(11.5), Inches(1),
         size=15, color=SUB, align=PP_ALIGN.CENTER, line_spacing=1.5)


# ──────────────────────────────────────────────────────────────────────────
# 04 — WHY NOW
# ──────────────────────────────────────────────────────────────────────────
s = add_slide(WHITE)
eyebrow(s, "Why now", Inches(0.8), Inches(0.8))
add_text(s, "Therapy works. But most fights happen between sessions.",
         Inches(0.8), Inches(1.4), Inches(12), Inches(1.4),
         size=36, bold=True, line_spacing=1.15)

cards = [
    ("Therapy is expensive",
     "$150 to $300 per session in most major cities. Most couples can't sustain weekly."),
    ("Waitlists are months long",
     "Couples in crisis wait 6 to 12 weeks for an opening. The fights don't wait."),
    ("The damage is between sessions",
     "Most ruptures happen at 11pm on a Tuesday, nine days from the next appointment."),
]
card_w = Inches(3.85)
gap = Inches(0.25)
start_left = (SW - (card_w * 3 + gap * 2)) // 2
top = Inches(3.4)
for i, (title, body) in enumerate(cards):
    left = start_left + (card_w + gap) * i
    add_card(s, left, top, card_w, Inches(2.6))
    add_text(s, title, left + Inches(0.3), top + Inches(0.35),
             card_w - Inches(0.6), Inches(0.7),
             size=18, bold=True, color=INK)
    add_text(s, body, left + Inches(0.3), top + Inches(1.2),
             card_w - Inches(0.6), Inches(2),
             size=13, color=SUB, line_spacing=1.5)

add_text(s,
         "AI has finally crossed the threshold to help people process emotion safely, "
         "privately, and in real time. The window to build this is open.",
         Inches(0.8), Inches(6.5), Inches(11.5), Inches(1),
         size=14, color=SUB, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────
# 05 — SOLUTION
# ──────────────────────────────────────────────────────────────────────────
s = add_slide(LIME_LIGHT)
eyebrow(s, "The solution", Inches(0.8), Inches(1.5), color=EYEBROW_GREEN)
add_text(s, "Hey Otis is a private guide that walks couples\nfrom rupture to repair.",
         Inches(0.5), Inches(2.4), Inches(12.3), Inches(2.5),
         size=46, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.15)
add_text(s,
         "In your pocket. At 11pm. When your therapist doesn't pick up.\n"
         "Not a replacement for therapy. The thing that makes therapy work.",
         Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.5),
         size=20, color=INK, align=PP_ALIGN.CENTER, line_spacing=1.5)


# ──────────────────────────────────────────────────────────────────────────
# 06 — HOW IT WORKS — 4 step cards 2x2
# ──────────────────────────────────────────────────────────────────────────
s = add_slide(CREAM)
eyebrow(s, "How it works", Inches(0.8), Inches(0.6))
add_text(s, "Four steps. One conversation at a time.",
         Inches(0.8), Inches(1.1), Inches(12), Inches(1),
         size=34, bold=True)

steps = [
    (1, TEAL, "Vent", "Let it all out",
     "A private space to say exactly what you're feeling. Putting emotion into words reduces its intensity."),
    (2, PERIWINKLE, "Understand", "What's really going on?",
     "Most arguments aren't about what they seem. Hey Otis helps you find the unmet need underneath."),
    (3, ORANGE, "Prepare", "Find the right words",
     "Coaches you to frame what you want to say, so your partner hears you, not an attack."),
    (4, PURPLE, "Nurture", "Have the conversation",
     "Step by step support during the actual repair. How to open softly, what to say if it gets heated."),
]
cw = Inches(5.85)
ch = Inches(2.4)
gap = Inches(0.3)
start_left = (SW - (cw * 2 + gap)) // 2
start_top = Inches(2.6)
for i, (num, col, name, sub, body) in enumerate(steps):
    row, colidx = divmod(i, 2)
    left = start_left + (cw + gap) * colidx
    top = start_top + (ch + gap) * row
    add_card(s, left, top, cw, ch)
    # number circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.35), top + Inches(0.35),
                                Inches(0.55), Inches(0.55))
    circle.fill.solid(); circle.fill.fore_color.rgb = col
    circle.line.fill.background()
    add_text(s, str(num), left + Inches(0.35), top + Inches(0.39),
             Inches(0.55), Inches(0.5), size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, line_spacing=1.0)
    add_text(s, name, left + Inches(1.05), top + Inches(0.32),
             cw - Inches(1.2), Inches(0.5), size=18, bold=True)
    add_text(s, sub.upper(), left + Inches(1.05), top + Inches(0.7),
             cw - Inches(1.2), Inches(0.4), size=10, bold=True, color=SUB)
    add_text(s, body, left + Inches(0.35), top + Inches(1.2),
             cw - Inches(0.7), Inches(1.2), size=13, color=BODY, line_spacing=1.5)


# ──────────────────────────────────────────────────────────────────────────
# 07 — PRODUCT (split: text left, image placeholder right)
# ──────────────────────────────────────────────────────────────────────────
s = add_slide(WHITE)
eyebrow(s, "The product", Inches(0.8), Inches(0.8))
add_text(s, "A relationship that learns the more you share.",
         Inches(0.8), Inches(1.4), Inches(7), Inches(1.6),
         size=34, bold=True, line_spacing=1.15)

bullets = [
    "Short assessments map how each partner experiences love, conflict, and stress",
    "Voice or text input. Use it solo or together",
    "End to end private. We never train models on your data. We never sell it.",
    "Available 24/7 on iOS, Android, and web",
]
ty = Inches(3.4)
for b in bullets:
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), ty + Inches(0.13),
                             Inches(0.18), Inches(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = LIME
    dot.line.fill.background()
    add_text(s, b, Inches(1.15), ty, Inches(6), Inches(0.6), size=15, color=BODY, line_spacing=1.5)
    ty += Inches(0.7)

# Try to embed the hero image — falls back to a card if missing.
try:
    s.shapes.add_picture(
        "/Users/samanthaahhee/tether/website/public/hero-image.png",
        Inches(7.7), Inches(1.6), width=Inches(5.0),
    )
except FileNotFoundError:
    add_card(s, Inches(7.7), Inches(1.6), Inches(5.0), Inches(5.0), fill=CREAM)
    add_text(s, "Product screen", Inches(7.7), Inches(3.8), Inches(5.0), Inches(0.5),
             size=14, color=SUB, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────
# 08 — FRAMEWORKS (5 pills)
# ──────────────────────────────────────────────────────────────────────────
s = add_slide(CREAM)
eyebrow(s, "Evidence based", Inches(0.8), Inches(0.8))
add_text(s, "Built on the same frameworks couples therapists use.",
         Inches(0.8), Inches(1.4), Inches(12), Inches(1.4),
         size=32, bold=True, line_spacing=1.2)

pills = [
    ("Gottman", "Conflict styles + sound\nrelationship house"),
    ("Attachment", "Anxious / avoidant /\nsecure"),
    ("NVC", "Observation, feeling,\nneed, request"),
    ("Love Languages", "How each partner\ngives + receives"),
    ("IFS", "Parts work for\nself understanding"),
]
pw = Inches(2.35)
gap = Inches(0.2)
total_w = pw * 5 + gap * 4
start_left = (SW - total_w) // 2
top = Inches(3.5)
for i, (t, sub) in enumerate(pills):
    left = start_left + (pw + gap) * i
    add_card(s, left, top, pw, Inches(2.0))
    add_text(s, t, left + Inches(0.15), top + Inches(0.4),
             pw - Inches(0.3), Inches(0.5), size=18, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, sub, left + Inches(0.15), top + Inches(1.0),
             pw - Inches(0.3), Inches(1.0), size=11, color=SUB,
             align=PP_ALIGN.CENTER, line_spacing=1.5)
add_text(s,
         "Otis is not making it up. Every prompt, every reframe, every step is grounded "
         "in peer reviewed couples research.",
         Inches(0.8), Inches(6.3), Inches(11.5), Inches(1),
         size=14, color=SUB, align=PP_ALIGN.CENTER, line_spacing=1.5)


# ──────────────────────────────────────────────────────────────────────────
# 09 — MARKET (3 horizontal rows)
# ──────────────────────────────────────────────────────────────────────────
s = add_slide(WHITE)
eyebrow(s, "The market", Inches(0.8), Inches(0.8))
add_text(s, "A massive, underserved category.",
         Inches(0.8), Inches(1.4), Inches(12), Inches(1),
         size=36, bold=True)

rows = [
    ("$1.2T", "TAM", "Global wellness market (Global Wellness Institute, 2024)"),
    ("$58B", "SAM", "Mental health + relationship wellness, growing 17% YoY"),
    ("$4B", "SOM (Year 5)", "English speaking couples 25 to 55 willing to pay for ongoing support"),
]
ry = Inches(2.7)
for size, lab, body in rows:
    add_card(s, Inches(0.8), ry, Inches(11.7), Inches(1.1))
    add_text(s, size, Inches(1.1), ry + Inches(0.18), Inches(2.5), Inches(0.8),
             size=32, bold=True, color=TEAL)
    add_text(s, lab, Inches(3.7), ry + Inches(0.2), Inches(3), Inches(0.4),
             size=13, bold=True)
    add_text(s, body, Inches(3.7), ry + Inches(0.55), Inches(8.5), Inches(0.5),
             size=12, color=SUB)
    ry += Inches(1.25)
add_text(s,
         "47M couples in the US alone are in long term relationships and not in therapy. "
         "We're not stealing share from BetterHelp. We're creating a new category.",
         Inches(0.8), Inches(6.6), Inches(11.5), Inches(1),
         size=13, color=SUB, align=PP_ALIGN.CENTER, line_spacing=1.5)


# ──────────────────────────────────────────────────────────────────────────
# 10 — BUSINESS MODEL (3 price cards)
# ──────────────────────────────────────────────────────────────────────────
s = add_slide(CREAM)
eyebrow(s, "Business model", Inches(0.8), Inches(0.8))
add_text(s, "Subscription. Aligned with the long term work of a relationship.",
         Inches(0.8), Inches(1.4), Inches(12), Inches(1.4),
         size=30, bold=True, line_spacing=1.15)

prices = [
    ("Free", "$0", "Daily check ins, one guided\nrepair per month, basic insights", False),
    ("Hey Otis Plus", "$14.99/mo", "Unlimited repairs, partner sync,\nfull assessment library", True),
    ("Couples", "$22/mo", "Both partners synced, shared\nhistory, deeper insight", False),
]
pw = Inches(3.85)
gap = Inches(0.25)
total_w = pw * 3 + gap * 2
start_left = (SW - total_w) // 2
top = Inches(3.3)
for i, (tier, price, body, feat) in enumerate(prices):
    left = start_left + (pw + gap) * i
    fill = WHITE
    add_card(s, left, top, pw, Inches(2.7), fill=fill,
             stroke=LIME if feat else RGBColor(0xDE, 0xDD, 0xE8))
    add_text(s, tier.upper(), left, top + Inches(0.35),
             pw, Inches(0.4), size=12, bold=True, color=SUB,
             align=PP_ALIGN.CENTER)
    add_text(s, price, left, top + Inches(0.95),
             pw, Inches(0.8), size=32, bold=True, color=INK,
             align=PP_ALIGN.CENTER)
    add_text(s, body, left + Inches(0.3), top + Inches(1.85),
             pw - Inches(0.6), Inches(1), size=12, color=SUB,
             align=PP_ALIGN.CENTER, line_spacing=1.5)

add_text(s,
         "Therapy is $200 a session. Hey Otis is less than two coffees a week, "
         "available the moment you need it.",
         Inches(0.8), Inches(6.5), Inches(11.5), Inches(1),
         size=14, color=SUB, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────
# 11 — TRACTION
# ──────────────────────────────────────────────────────────────────────────
s = add_slide(WHITE)
eyebrow(s, "Traction", Inches(0.8), Inches(0.8))
add_text(s, "Where we are today.",
         Inches(0.8), Inches(1.4), Inches(12), Inches(1),
         size=36, bold=True)

big = [
    ("[X]", "Waitlist signups", "from heyotis.app + organic Reddit + IG"),
    ("[X]", "Survey responses", "primary research from real couples"),
    ("78%", "Said they'd try it", "from in-survey product validation"),
]
cw = Inches(3.85)
gap = Inches(0.25)
start_left = (SW - (cw * 3 + gap * 2)) // 2
top = Inches(3.0)
for i, (n, lab, sub) in enumerate(big):
    left = start_left + (cw + gap) * i
    add_card(s, left, top, cw, Inches(2.7))
    add_text(s, n, left, top + Inches(0.5),
             cw, Inches(1), size=56, bold=True, color=INK,
             align=PP_ALIGN.CENTER, line_spacing=1.0)
    add_text(s, lab.upper(), left, top + Inches(1.7),
             cw, Inches(0.4), size=12, bold=True, color=INK,
             align=PP_ALIGN.CENTER)
    add_text(s, sub, left + Inches(0.3), top + Inches(2.15),
             cw - Inches(0.6), Inches(0.6), size=11, color=SUB,
             align=PP_ALIGN.CENTER)

add_text(s,
         "Beta launching Q3 2026 in EN speaking markets. Apple + Google submission ready. "
         "Brand identity, marketing site, and waitlist infrastructure all live.",
         Inches(0.8), Inches(6.4), Inches(11.5), Inches(1),
         size=14, color=SUB, align=PP_ALIGN.CENTER, line_spacing=1.5)


# ──────────────────────────────────────────────────────────────────────────
# 12 — TEAM
# ──────────────────────────────────────────────────────────────────────────
s = add_slide(CREAM)
eyebrow(s, "Team", Inches(0.8), Inches(0.8))
add_text(s, "Built by people who've lived this.",
         Inches(0.8), Inches(1.4), Inches(12), Inches(1),
         size=36, bold=True)

team = [
    ("Samantha Ahhee", "Founder + CEO",
     "Designer + builder. Previously [your prior role]. Built Hey Otis after experiencing first hand what was missing between therapy and the everyday."),
    ("[Add advisor name]", "Clinical Advisor",
     "Licensed couples therapist. Validates the frameworks, the prompts, the safety guardrails."),
    ("[Add advisor name]", "Technical Advisor",
     "ML / AI background. Helps shape the model layer + private architecture."),
]
cw = Inches(3.85)
gap = Inches(0.25)
start_left = (SW - (cw * 3 + gap * 2)) // 2
top = Inches(2.8)
for i, (name, role, body) in enumerate(team):
    left = start_left + (cw + gap) * i
    add_card(s, left, top, cw, Inches(3.6))
    add_text(s, name, left + Inches(0.3), top + Inches(0.35),
             cw - Inches(0.6), Inches(0.6), size=18, bold=True)
    add_text(s, role.upper(), left + Inches(0.3), top + Inches(0.95),
             cw - Inches(0.6), Inches(0.4), size=11, bold=True, color=EYEBROW_GREEN)
    add_text(s, body, left + Inches(0.3), top + Inches(1.5),
             cw - Inches(0.6), Inches(2), size=13, color=BODY, line_spacing=1.5)


# ──────────────────────────────────────────────────────────────────────────
# 13 — ASK
# ──────────────────────────────────────────────────────────────────────────
s = add_slide(LIME_LIGHT)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.5), SW, Inches(3))
band.line.fill.background()
band.fill.solid(); band.fill.fore_color.rgb = LIME
eyebrow(s, "The ask", Inches(0.8), Inches(1.4))
add_text(s, "Raising [€XXX]K pre seed",
         Inches(0.5), Inches(2.2), Inches(12.3), Inches(2),
         size=64, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.1)
add_text(s,
         "To launch publicly, hit 10,000 paying users in 18 months,\n"
         "and prove that couples will pay to learn how to repair.",
         Inches(0.5), Inches(4.6), Inches(12.3), Inches(1.5),
         size=20, align=PP_ALIGN.CENTER, line_spacing=1.5)
add_text(s, "samantha.ahhee@gmail.com",
         Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
         size=18, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "heyotis.app",
         Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
         size=15, color=INK, align=PP_ALIGN.CENTER)


# ── Save ────────────────────────────────────────────────────────────────
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "HeyOtis_Pitch.pptx")
prs.save(out)
print(f"\n✓ Wrote {out}")
print(f"  Slides: {len(prs.slides)}")
